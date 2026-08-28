"""Prueba de humo del convertidor, sin abrir la ventana.

Genera archivos de muestra de varios formatos, los convierte con `app.converter`
y reporta cuántos caracteres salieron. Incluye a propósito un PDF corrupto para
comprobar que el error se traduce a un mensaje legible en vez de tumbar el proceso.

Uso: uv run python scripts/smoke_convert.py
"""

from __future__ import annotations

import json
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from app.converter import ConversionError, convert  # noqa: E402


def minimal_pdf(text: str) -> bytes:
    """Arma un PDF válido mínimo con una línea de texto.

    Evita depender de una librería de escritura de PDF solo para la prueba.
    La tabla xref se construye con los desplazamientos reales, así pdfminer lo
    lee sin reconstruirla.
    """
    stream = f"BT /F1 14 Tf 24 120 Td ({text}) Tj ET".encode("latin-1")
    objects = [
        b"<</Type/Catalog/Pages 2 0 R>>",
        b"<</Type/Pages/Kids[3 0 R]/Count 1>>",
        b"<</Type/Page/Parent 2 0 R/MediaBox[0 0 200 200]"
        b"/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>",
        b"<</Length " + str(len(stream)).encode() + b">>stream\n" + stream + b"\nendstream",
        b"<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>",
    ]

    out = bytearray(b"%PDF-1.4\n")
    offsets: list[int] = []
    for index, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{index} 0 obj\n".encode() + body + b"\nendobj\n"

    xref_at = len(out)
    out += f"xref\n0 {len(objects) + 1}\n".encode()
    out += b"0000000000 65535 f \n"
    for offset in offsets:
        out += f"{offset:010d} 00000 n \n".encode()
    out += (
        f"trailer<</Size {len(objects) + 1}/Root 1 0 R>>\n"
        f"startxref\n{xref_at}\n%%EOF\n"
    ).encode()

    return bytes(out)


def build_samples(folder: Path) -> list[Path]:
    """Escribe un archivo de muestra por formato y devuelve las rutas."""
    samples: list[Path] = []

    txt = folder / "notas.txt"
    txt.write_text("Primera línea\nSegunda línea con acentos: ñ á é\n", encoding="utf-8")
    samples.append(txt)

    md = folder / "guía.md"
    md.write_text("# Título\n\nUn párrafo con **negrita**.\n", encoding="utf-8")
    samples.append(md)

    csv = folder / "datos.csv"
    csv.write_text("nombre,cantidad\nlápiz,3\ncuaderno,7\n", encoding="utf-8")
    samples.append(csv)

    js = folder / "config.json"
    js.write_text(json.dumps({"nombre": "ToMarkdown", "versión": 1}, ensure_ascii=False))
    samples.append(js)

    html = folder / "página.html"
    html.write_text(
        "<html><body><h1>Encabezado</h1><p>Texto <b>marcado</b>.</p></body></html>",
        encoding="utf-8",
    )
    samples.append(html)

    xml = folder / "feed.xml"
    xml.write_text("<root><item>uno</item><item>dos</item></root>", encoding="utf-8")
    samples.append(xml)

    notebook = folder / "análisis.ipynb"
    notebook.write_text(
        json.dumps(
            {
                "cells": [
                    {"cell_type": "markdown", "metadata": {}, "source": ["# Cuaderno\n"]},
                    {
                        "cell_type": "code",
                        "metadata": {},
                        "execution_count": 1,
                        "outputs": [],
                        "source": ["print('hola')\n"],
                    },
                ],
                "metadata": {},
                "nbformat": 4,
                "nbformat_minor": 5,
            }
        ),
        encoding="utf-8",
    )
    samples.append(notebook)

    try:
        from openpyxl import Workbook

        book = Workbook()
        sheet = book.active
        sheet.append(["producto", "stock"])
        sheet.append(["cuaderno", 12])
        xlsx = folder / "inventario.xlsx"
        book.save(xlsx)
        samples.append(xlsx)
    except ImportError:
        print("  (openpyxl no disponible, se omite xlsx)")

    try:
        from pptx import Presentation

        deck = Presentation()
        slide = deck.slides.add_slide(deck.slide_layouts[0])
        slide.shapes.title.text = "Diapositiva de prueba"
        pptx = folder / "presentación.pptx"
        deck.save(pptx)
        samples.append(pptx)
    except ImportError:
        print("  (python-pptx no disponible, se omite pptx)")

    pdf_bytes = minimal_pdf("Informe anual de prueba")
    pdf = folder / "informe.pdf"
    pdf.write_bytes(pdf_bytes)
    samples.append(pdf)

    # PDF deliberadamente roto: el mismo, cortado a la mitad. La cola tiene que
    # marcarlo como error y seguir con el resto.
    broken = folder / "roto.pdf"
    broken.write_bytes(pdf_bytes[: len(pdf_bytes) // 2])
    samples.append(broken)

    return samples


def main() -> int:
    """Convierte las muestras e informa el resultado. Devuelve el código de salida."""
    with tempfile.TemporaryDirectory() as tmp:
        folder = Path(tmp)
        samples = build_samples(folder)

        print(f"\nConvirtiendo {len(samples)} archivos de muestra\n")
        failures = 0

        for sample in samples:
            expected_failure = sample.name == "roto.pdf"
            try:
                markdown = convert(str(sample))
            except ConversionError as exc:
                mark = "✓" if expected_failure else "✗"
                if not expected_failure:
                    failures += 1
                print(f"  {mark} {sample.name:<24} error: {exc}")
            else:
                if expected_failure:
                    print(f"  ! {sample.name:<24} se esperaba un error y convirtió")
                else:
                    print(f"  ✓ {sample.name:<24} {len(markdown):>6} caracteres")

        print()
        if failures:
            print(f"🔴 {failures} formatos fallaron")
            return 1

        print("✅ Todos los formatos esperados convirtieron")
        return 0


if __name__ == "__main__":
    raise SystemExit(main())
