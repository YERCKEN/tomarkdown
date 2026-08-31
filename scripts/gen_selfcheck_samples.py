"""Genera un archivo minimo de cada formato pesado para `--self-check`.

`--self-check` sin argumentos solo convierte una muestra de texto: no ejercita
los converters que dependen de un `hiddenimport` frágil (`pdf`, `docx`, `xlsx`,
`xls`, `pptx`). Este script escribe una muestra valida y diminuta de cada uno en
una carpeta, para que CI corra `--self-check <carpeta>` sobre el binario ya
empaquetado y detecte si falta una dependencia en el bundle.

No se versionan: se generan en cada corrida. Solo stdlib salvo `openpyxl` /
`python-pptx` (deps de runtime) y `xlwt` (dep de dev).

Uso:

    uv run python scripts/gen_selfcheck_samples.py <carpeta_destino>

`tests/conftest.py` reutiliza `minimal_pdf` desde acá.
"""

from __future__ import annotations

import io
import sys
import zipfile
from pathlib import Path

#: Texto que cada muestra incluye, para poder afirmar que la conversion trajo
#: contenido y no una cadena vacia.
MARKER = "Muestra de ToMarkdown para --self-check"

#: Muestras de texto plano: extension -> contenido. Van de yapa; el nucleo de
#: markitdown no las puede romper por un hiddenimport faltante, pero cuestan cero.
TEXT_SAMPLES: dict[str, str] = {
    "txt": f"{MARKER}\nSegunda linea con acentos: ñ á é\n",
    "md": f"# {MARKER}\n\nUn parrafo con **negrita**.\n",
    "csv": f"campo,valor\nmarcador,{MARKER}\ncantidad,3\n",
    "json": '{"marcador": "' + MARKER + '", "version": 1}\n',
    "xml": f"<root><item>{MARKER}</item><item>dos</item></root>\n",
    "html": f"<html><body><h1>{MARKER}</h1><p>Texto <b>marcado</b>.</p></body></html>\n",
}


def minimal_pdf(text: str) -> bytes:
    """Arma un PDF valido minimo con una linea de texto.

    Evita depender de una libreria de escritura de PDF. La tabla xref lleva los
    desplazamientos reales, asi pdfminer lo lee sin reconstruirla.
    """
    stream = f"BT /F1 14 Tf 24 120 Td ({text}) Tj ET".encode("latin-1")
    objects = [
        b"<</Type/Catalog/Pages 2 0 R>>",
        b"<</Type/Pages/Kids[3 0 R]/Count 1>>",
        b"<</Type/Page/Parent 2 0 R/MediaBox[0 0 200 200]"
        b"/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>",
        b"<</Length " + str(len(stream)).encode() + b">>stream\n" + stream + b"\nendstream",
        b"<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>",
    ]

    out = bytearray(b"%PDF-1.4\n")
    offsets: list[int] = []
    for index, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{index} 0 obj\n".encode() + body + b"\nendobj\n"

    xref_at = len(out)
    out += f"xref\n0 {len(objects) + 1}\n".encode()
    out += b"0000000000 65535 f \n"
    for offset in offsets:
        out += f"{offset:010d} 00000 n \n".encode()
    out += (
        f"trailer<</Size {len(objects) + 1}/Root 1 0 R>>\n"
        f"startxref\n{xref_at}\n%%EOF\n"
    ).encode()

    return bytes(out)


def minimal_docx(text: str) -> bytes:
    """Arma un .docx valido minimo (OOXML) sin depender de python-docx.

    Un .docx es un zip con tres partes: el mapa de tipos, la relacion raiz y el
    documento. mammoth (lo que usa markitdown) lo lee sin problema.
    """
    content_types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/word/document.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
        "</Types>"
    )
    rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
        'Target="word/document.xml"/>'
        "</Relationships>"
    )
    document = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        f"<w:body><w:p><w:r><w:t>{text}</w:t></w:r></w:p></w:body>"
        "</w:document>"
    )

    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", content_types)
        archive.writestr("_rels/.rels", rels)
        archive.writestr("word/document.xml", document)
    return buffer.getvalue()


def minimal_xlsx(text: str) -> bytes:
    """Arma un .xlsx minimo con openpyxl (dep de runtime)."""
    import openpyxl

    book = openpyxl.Workbook()
    sheet = book.active
    sheet.append(["marcador", "cantidad"])
    sheet.append([text, 3])
    buffer = io.BytesIO()
    book.save(buffer)
    return buffer.getvalue()


def minimal_xls(text: str) -> bytes:
    """Arma un .xls (BIFF) minimo con xlwt (dep de dev)."""
    import xlwt

    book = xlwt.Workbook()
    sheet = book.add_sheet("Hoja1")
    sheet.write(0, 0, "marcador")
    sheet.write(0, 1, "cantidad")
    sheet.write(1, 0, text)
    sheet.write(1, 1, 3)
    buffer = io.BytesIO()
    book.save(buffer)
    return buffer.getvalue()


def minimal_pptx(text: str) -> bytes:
    """Arma un .pptx minimo con python-pptx (dep de runtime)."""
    import pptx

    deck = pptx.Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[0])
    slide.shapes.title.text = text
    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def write_samples(destination: Path) -> list[Path]:
    """Escribe una muestra de cada formato en `destination`. Devuelve las rutas."""
    destination.mkdir(parents=True, exist_ok=True)
    written: list[Path] = []

    for extension, content in TEXT_SAMPLES.items():
        path = destination / f"muestra.{extension}"
        path.write_text(content, encoding="utf-8")
        written.append(path)

    binary = {
        "pdf": minimal_pdf(MARKER),
        "docx": minimal_docx(MARKER),
        "xlsx": minimal_xlsx(MARKER),
        "xls": minimal_xls(MARKER),
        "pptx": minimal_pptx(MARKER),
    }
    for extension, data in binary.items():
        path = destination / f"muestra.{extension}"
        path.write_bytes(data)
        written.append(path)

    # Un zip con la muestra de texto adentro, para ejercitar el ZipConverter.
    zip_path = destination / "muestra.zip"
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("dentro.txt", TEXT_SAMPLES["txt"])
    written.append(zip_path)

    return written


def main(argv: list[str]) -> int:
    if len(argv) != 1:
        print("uso: gen_selfcheck_samples.py <carpeta_destino>", file=sys.stderr)
        return 2

    destination = Path(argv[0])
    written = write_samples(destination)
    print(f"✅ {len(written)} muestras en {destination}")
    for path in written:
        print(f"  · {path.name} ({path.stat().st_size} bytes)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
